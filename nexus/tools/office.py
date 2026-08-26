"""Deterministic Office deliverables — PPTX / PDF / DOCX (no LLM guessing bytes)."""
from __future__ import annotations

import importlib
import subprocess
import sys
from pathlib import Path
from typing import Any, List, Optional

from .base import Risk, ToolRegistry, ToolResult
from .paths import in_workspace


def _ensure(mod: str, pip_name: Optional[str] = None) -> Any:
    try:
        return importlib.import_module(mod)
    except ImportError:
        pkg = pip_name or mod
        subprocess.run(
            [sys.executable, "-m", "pip", "install", "--quiet", pkg],
            check=False, capture_output=True, timeout=180,
        )
        return importlib.import_module(mod)


class OfficeTools:
    def __init__(self, workspace: Path) -> None:
        self.root = Path(workspace).resolve()

    def _out(self, path: str) -> Path:
        p = Path(path)
        if not p.is_absolute():
            p = self.root / p
        p = p.resolve()
        if not in_workspace(p, self.root):
            raise ValueError("path escapes workspace")
        p.parent.mkdir(parents=True, exist_ok=True)
        return p

    def make_pptx(self, path: str, title: str, slides: str) -> ToolResult:
        """slides: '##' heading then body lines, slides separated by a line with only ---"""
        try:
            pptx = _ensure("pptx", "python-pptx")
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except Exception as e:
            return ToolResult(False, error=f"python-pptx unavailable: {e}")
        try:
            dest = self._out(path if path.endswith(".pptx") else path + ".pptx")
            prs = pptx.Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blocks = [b.strip() for b in (slides or "").split("\n---\n") if b.strip()]
            if not blocks:
                blocks = [f"## {title}\nOverview"]
            for i, block in enumerate(blocks):
                lines = [ln.rstrip() for ln in block.splitlines() if ln.strip()]
                head = lines[0].lstrip("#").strip() if lines else f"Slide {i+1}"
                body = lines[1:] if len(lines) > 1 else []
                layout = prs.slide_layouts[6]  # blank
                sl = prs.slides.add_slide(layout)
                box = sl.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(1.1))
                tf = box.text_frame
                tf.text = head
                tf.paragraphs[0].font.size = Pt(32)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
                if i == 0 and title and title not in head:
                    sub = sl.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
                    sub.text_frame.text = title
                    sub.text_frame.paragraphs[0].font.size = Pt(16)
                tb = sl.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.8), Inches(4.8))
                tff = tb.text_frame
                tff.word_wrap = True
                for j, ln in enumerate(body[:12]):
                    p = tff.paragraphs[0] if j == 0 else tff.add_paragraph()
                    p.text = ln.lstrip("-• ").strip()
                    p.level = 0
                    p.font.size = Pt(20)
                    p.space_after = Pt(8)
            prs.save(str(dest))
            return ToolResult(True, output=f"PPTX written: {dest.relative_to(self.root)} ({dest.stat().st_size} bytes, {len(prs.slides)} slides)")
        except Exception as e:  # noqa: BLE001
            return ToolResult(False, error=f"make_pptx failed: {e}")

    def make_pdf(self, path: str, title: str, body: str) -> ToolResult:
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.enums import TA_LEFT
        except ImportError:
            try:
                _ensure("reportlab")
                from reportlab.lib.pagesizes import A4
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.lib.units import mm
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.enums import TA_LEFT
            except Exception as e:
                return ToolResult(False, error=f"reportlab unavailable: {e}")
        try:
            dest = self._out(path if path.endswith(".pdf") else path + ".pdf")
            doc = SimpleDocTemplate(str(dest), pagesize=A4,
                                    leftMargin=18 * mm, rightMargin=18 * mm,
                                    topMargin=16 * mm, bottomMargin=16 * mm)
            styles = getSampleStyleSheet()
            styles.add(ParagraphStyle(name="HNexus", parent=styles["Heading1"], fontSize=18, spaceAfter=10))
            styles.add(ParagraphStyle(name="BNexus", parent=styles["BodyText"], fontSize=11,
                                      leading=15, alignment=TA_LEFT, spaceAfter=6))
            story: List = [Paragraph(_esc(title or "Document"), styles["HNexus"]), Spacer(1, 8)]
            for para in (body or "").split("\n"):
                if not para.strip():
                    story.append(Spacer(1, 6))
                    continue
                story.append(Paragraph(_esc(para), styles["BNexus"]))
            doc.build(story)
            return ToolResult(True, output=f"PDF written: {dest.relative_to(self.root)} ({dest.stat().st_size} bytes)")
        except Exception as e:  # noqa: BLE001
            return ToolResult(False, error=f"make_pdf failed: {e}")

    def make_docx(self, path: str, title: str, body: str) -> ToolResult:
        try:
            docx = _ensure("docx", "python-docx")
        except Exception as e:
            return ToolResult(False, error=f"python-docx unavailable: {e}")
        try:
            dest = self._out(path if path.endswith(".docx") else path + ".docx")
            d = docx.Document()
            d.add_heading(title or "Document", level=1)
            for para in (body or "").split("\n"):
                if para.startswith("## "):
                    d.add_heading(para[3:].strip(), level=2)
                elif para.startswith("# "):
                    d.add_heading(para[2:].strip(), level=1)
                else:
                    d.add_paragraph(para)
            d.save(str(dest))
            return ToolResult(True, output=f"DOCX written: {dest.relative_to(self.root)} ({dest.stat().st_size} bytes)")
        except Exception as e:  # noqa: BLE001
            return ToolResult(False, error=f"make_docx failed: {e}")

    def register(self, reg: ToolRegistry) -> None:
        S = {"type": "string"}
        who = ["coder", "worker", "supervisor", "solo"]
        reg.add("make_pptx",
                "Create a real .pptx PowerPoint. slides = markdown blocks separated by a line with only --- ; each block starts with ## title.",
                {"type": "object", "properties": {"path": S, "title": S, "slides": S},
                 "required": ["path", "slides"]},
                self.make_pptx, Risk.WRITE, agents=who)
        reg.add("make_pdf",
                "Create a real .pdf document from title + body text (paragraphs separated by newlines).",
                {"type": "object", "properties": {"path": S, "title": S, "body": S},
                 "required": ["path", "body"]},
                self.make_pdf, Risk.WRITE, agents=who)
        reg.add("make_docx",
                "Create a real .docx Word document. body may use # / ## headings.",
                {"type": "object", "properties": {"path": S, "title": S, "body": S},
                 "required": ["path", "body"]},
                self.make_docx, Risk.WRITE, agents=who)


def _esc(s: str) -> str:
    return (s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
